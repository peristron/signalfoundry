#  SIGNAL FOUNDRY (an UN-structured data intel engine)
#  prod ready (v2.9 - robust graphics safety)
#  arhitecture: hybrid streaming + "data refinery" utility
#
import io
import os
import re
import html
import gc
import time
import csv
import json
import math
import string
import zipfile
import tempfile
import logging
import secrets
from dataclasses import dataclass, field
from urllib.parse import urlparse
from collections import Counter, defaultdict
from typing import Dict, List, Tuple, Iterable, Optional, Callable, Any, Union, Set
from datetime import datetime

import numpy as np
import pandas as pd
import streamlit as st
import matplotlib.pyplot as plt
import matplotlib.dates as mdates
from io import BytesIO
from wordcloud import WordCloud, STOPWORDS
from matplotlib import font_manager
from itertools import pairwise
import openai

# --- graph imports
import networkx as nx
import networkx.algorithms.community as nx_comm
from streamlit_agraph import agraph, Node, Edge, Config

# -3rd party imports checks
try:
    import requests
    from bs4 import BeautifulSoup
except ImportError:
    requests = None
    BeautifulSoup = None

try:
    import qrcode
except ImportError:
    qrcode = None

try:
    from scipy.stats import beta as beta_dist
except ImportError:
    beta_dist = None

try:
    from sklearn.decomposition import LatentDirichletAllocation, NMF
    from sklearn.feature_extraction import DictVectorizer
except ImportError:
    LatentDirichletAllocation = None
    NMF = None
    DictVectorizer = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    from dateutil import parser as date_parser
except ImportError:
    date_parser = None

try:
    import pypdf
except ImportError:
    pypdf = None

try:
    import pptx
except ImportError:
    pptx = None

try:
    import nltk
    from nltk.sentiment.vader import SentimentIntensityAnalyzer
    from nltk.stem import WordNetLemmatizer
except ImportError:
    nltk = None
    SentimentIntensityAnalyzer = None
    WordNetLemmatizer = None


# ⚙️ constants and config
# ========================================

MAX_TOPIC_DOCS = 50_000
MAX_SPEAKER_NAME_LENGTH = 30
SENTIMENT_ANALYSIS_TOP_N = 5000
URL_SCRAPE_RATE_LIMIT_SECONDS = 1.0
PROGRESS_UPDATE_MIN_INTERVAL = 100
NPMI_MIN_FREQ = 3
MAX_FILE_SIZE_MB = 200

# regex patterns
HTML_TAG_RE = re.compile(r"<[^>]+>")
CHAT_ARTIFACT_RE = re.compile(
    r":\w+:"
    r"|\b(?:monday|tuesday|wednesday|thursday|friday|saturday|sunday|today|yesterday) at \d{1,2}:\d{2}\b"
    r"|\b\d+\s+repl(?:y|ies)\b"
    r"|\d{2}:\d{2}:\d{2}\.\d{3}\s+-->\s+\d{2}:\d{2}:\d{2}\.\d{3}"
    r"|\[[^\]]+\]",
    flags=re.IGNORECASE
)
URL_EMAIL_RE = re.compile(
    r'(?:https?://|www\.)[a-zA-Z0-9-]+(?:\.[a-zA-Z0-9-]+)+[^\s]*'
    r'|(?:[a-zA-Z0-9._%+-]+@[a-zA-Z0-9.-]+\.[a-zA-Z]{2,})',
    flags=re.IGNORECASE
)
# NER lite pattern (capitalized words in sequence)
NER_CAPS_RE = re.compile(r'\b([A-Z][a-z]+(?:\s+[A-Z][a-z]+)*)\b')

# logger setup
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger("IntelEngine")

# custom exceptions
class ReaderError(Exception):
    pass

class ValidationError(Exception):
    pass

# 📦 dataclassed
# ===========================================

@dataclass
class CleaningConfig:
    remove_chat: bool = True
    remove_html: bool = True
    remove_urls: bool = True
    unescape: bool = True
    phrase_pattern: Optional[re.Pattern] = None

@dataclass
class ProcessingConfig:
    min_word_len: int = 2
    drop_integers: bool = True
    compute_bigrams: bool = True
    use_lemmatization: bool = False
    translate_map: Dict[int, Optional[int]] = field(default_factory=dict)
    stopwords: Set[str] = field(default_factory=set)


# 🛡️ security and validation utils
# ==========================================

def get_auth_password() -> str:
    pwd = st.secrets.get("auth_password")
    if not pwd:
        st.error("🚨 Configuration Error: 'auth_password' not set in .streamlit/secrets.toml.")
        st.stop()
    return pwd

def validate_url(url: str) -> bool:
    try:
        parsed = urlparse(url)
        if parsed.scheme not in ('http', 'https'):
            return False
        if parsed.hostname in ('localhost', '127.0.0.1', '0.0.0.0', '::1'):
            return False
        return True
    except Exception:
        return False

def validate_sketch_data(data: Dict) -> bool:
    REQUIRED_KEYS = {"total_rows", "counts", "bigrams", "topic_docs"}
    if not isinstance(data, dict): return False
    if not REQUIRED_KEYS.issubset(data.keys()): return False
    return True


# 🧠 core logic, scanner
# =========================================

class StreamScanner:
    def __init__(self, doc_batch_size=5):
        self.global_counts = Counter()
        self.global_bigrams = Counter()
        self.total_rows_processed = 0
        self.topic_docs: List[Counter] = []
        
        # newer analytical structures (v2.0+)
        self.temporal_counts = defaultdict(Counter) # { '2023-10-01': {'word': 5} }
        self.category_counts = defaultdict(Counter) # { 'CategoryA': {'word': 10} }
        self.doc_freqs = Counter() # DF for TF-IDF
        self.entity_counts = Counter() # NER lite storage
        
        self.DOC_BATCH_SIZE = doc_batch_size
        self.limit_reached = False

    def set_batch_size(self, size: int):
        self.DOC_BATCH_SIZE = size

    def update_global_stats(self, counts: Counter, bigrams: Counter, rows: int):
        self.global_counts.update(counts)
        self.global_bigrams.update(bigrams)
        self.total_rows_processed += rows

    def add_topic_sample(self, doc_counts: Counter):
        if not doc_counts: return
        # update document frequency (presence check) for TF-IDF
        self.doc_freqs.update(doc_counts.keys())
        
        if self.limit_reached: return
        if len(self.topic_docs) >= MAX_TOPIC_DOCS:
            self.limit_reached = True
            return
        self.topic_docs.append(doc_counts)
    
    def update_metadata_stats(self, date_key: Optional[str], cat_key: Optional[str], tokens: List[str]):
        if date_key:
            self.temporal_counts[date_key].update(tokens)
        if cat_key:
            self.category_counts[cat_key].update(tokens)

    def update_entities(self, entities: List[str]):
        if entities:
            self.entity_counts.update(entities)

    def to_json(self) -> str:
        # simplifying complex structures for JSON serialization
        serializable_bigrams = {f"{k[0]}|{k[1]}": v for k, v in self.global_bigrams.items()}
        data = {
            "total_rows": self.total_rows_processed,
            "counts": dict(self.global_counts),
            "bigrams": serializable_bigrams,
            "topic_docs": [dict(c) for c in self.topic_docs],
            "limit_reached": self.limit_reached,
            # Persistence for new features
            "temporal_counts": {k: dict(v) for k, v in self.temporal_counts.items()},
            "entity_counts": dict(self.entity_counts),
            "doc_freqs": dict(self.doc_freqs)
        }
        return json.dumps(data)

    def load_from_json(self, json_str: str) -> bool:
        try:
            data = json.loads(json_str)
            if not validate_sketch_data(data): return False
            
            self.total_rows_processed = data.get("total_rows", 0)
            self.global_counts = Counter(data.get("counts", {}))
            
            raw_bigrams = data.get("bigrams", {})
            self.global_bigrams = Counter()
            for k, v in raw_bigrams.items():
                if "|" in k:
                    parts = k.split("|", 1)
                    self.global_bigrams[(parts[0], parts[1])] = v
            
            self.topic_docs = [Counter(d) for d in data.get("topic_docs", [])]
            self.limit_reached = data.get("limit_reached", False)
            
            # new
            self.entity_counts = Counter(data.get("entity_counts", {}))
            self.doc_freqs = Counter(data.get("doc_freqs", {}))
            
            raw_temp = data.get("temporal_counts", {})
            self.temporal_counts = defaultdict(Counter)
            for k, v in raw_temp.items():
                self.temporal_counts[k] = Counter(v)
                
            return True
        except Exception as e:
            logger.error(f"JSON Load Error: {e}")
            return False

# session state init
if 'sketch' not in st.session_state: st.session_state['sketch'] = StreamScanner()
if 'total_cost' not in st.session_state: st.session_state['total_cost'] = 0.0
if 'total_tokens' not in st.session_state: st.session_state['total_tokens'] = 0
if 'authenticated' not in st.session_state: st.session_state['authenticated'] = False
if 'auth_error' not in st.session_state: st.session_state['auth_error'] = False
if 'ai_response' not in st.session_state: st.session_state['ai_response'] = ""
if 'last_sketch_hash' not in st.session_state: st.session_state['last_sketch_hash'] = None

def reset_sketch():
    st.session_state['sketch'] = StreamScanner()
    st.session_state['ai_response'] = ""
    st.session_state['last_sketch_hash'] = None
    gc.collect()

def perform_login():
    try:
        correct_password = get_auth_password()
        if secrets.compare_digest(st.session_state.password_input, correct_password):
            st.session_state['authenticated'] = True
            st.session_state['auth_error'] = False
            st.session_state['password_input'] = ""
        else:
            st.session_state['auth_error'] = True
    except Exception:
        st.session_state['auth_error'] = True

def logout():
    st.session_state['authenticated'] = False
    st.session_state['ai_response'] = ""


# 🛠️ helpers, setup
# ============================================

@st.cache_resource(show_spinner="Init NLTK...")
def setup_nlp_resources():
    if nltk is None: return None, None
    try: 
        nltk.data.find('sentiment/vader_lexicon.zip')
        nltk.data.find('corpora/wordnet.zip')
        nltk.data.find('corpora/omw-1.4.zip')
    except LookupError: 
        try:
            nltk.download('vader_lexicon')
            nltk.download('wordnet')
            nltk.download('omw-1.4')
        except:
            pass
    
    sia = SentimentIntensityAnalyzer()
    lemmatizer = WordNetLemmatizer()
    return sia, lemmatizer

@st.cache_data(show_spinner=False)
def list_system_fonts() -> Dict[str, str]:
    mapping = {}
    for fe in font_manager.fontManager.ttflist:
        if fe.name not in mapping: mapping[fe.name] = fe.fname
    return dict(sorted(mapping.items(), key=lambda x: x[0].lower()))

def build_punct_translation(keep_hyphens: bool, keep_apostrophes: bool) -> dict:
    punct = string.punctuation
    if keep_hyphens: punct = punct.replace("-", "")
    if keep_apostrophes: punct = punct.replace("'", "")
    return str.maketrans("", "", punct)

def parse_user_stopwords(raw: str) -> Tuple[List[str], List[str]]:
    raw = raw.replace("\n", ",").replace(".", ",")
    phrases, singles = [], []
    for item in [x.strip() for x in raw.split(",") if x.strip()]:
        if " " in item: phrases.append(item.lower())
        else: singles.append(item.lower())
    return phrases, singles

def default_prepositions() -> set:
    return {'about', 'above', 'across', 'after', 'against', 'along', 'among', 'around', 'at', 'before', 'behind', 'below', 'beneath', 'beside', 'between', 'beyond', 'but', 'by', 'concerning', 'despite', 'down', 'during', 'except', 'for', 'from', 'in', 'inside', 'into', 'like', 'near', 'of', 'off', 'on', 'onto', 'out', 'outside', 'over', 'past', 'regarding', 'since', 'through', 'throughout', 'to', 'toward', 'under', 'underneath', 'until', 'up', 'upon', 'with', 'within', 'without'}

def build_phrase_pattern(phrases: List[str]) -> Optional[re.Pattern]:
    if not phrases: return None
    escaped = [re.escape(p) for p in phrases if p]
    if not escaped: return None
    return re.compile(rf"\b(?:{'|'.join(escaped)})\b", flags=re.IGNORECASE)

def estimate_row_count_from_bytes(file_bytes: bytes) -> int:
    if not file_bytes: return 0
    return file_bytes.count(b'\n') + 1

def make_unique_header(raw_names: List[Optional[str]]) -> List[str]:
    seen: Dict[str, int] = {}
    result: List[str] = []
    for i, nm in enumerate(raw_names):
        name = (str(nm).strip() if nm is not None else "")
        if not name: name = f"col_{i}"
        if name in seen:
            seen[name] += 1
            unique = f"{name}__{seen[name]}"
        else:
            seen[name] = 1
            unique = name
        result.append(unique)
    return result

def extract_entities_regex(text: str, stopwords: Set[str]) -> List[str]:
    # lightweight NER without heavy models
    candidates = NER_CAPS_RE.findall(text)
    valid = []
    for c in candidates:
        # filter out if it's just a common stopword capitalized at start of sentence
        if c.lower() in stopwords: continue
        if len(c) < 3: continue
        valid.append(c)
    return valid

# --- virtual files and web/url
class VirtualFile:
    def __init__(self, name: str, text_content: str):
        self.name = name
        self._bytes = text_content.encode('utf-8')
    
    def getvalue(self) -> bytes:
        return self._bytes
    
    def getbuffer(self) -> memoryview:
        return memoryview(self._bytes)

def fetch_url_content(url: str) -> Optional[str]:
    if not requests or not BeautifulSoup: return None
    if not validate_url(url): return None
    try:
        headers = {'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36'}
        response = requests.get(url, headers=headers, timeout=10)
        response.raise_for_status()
        soup = BeautifulSoup(response.content, 'html.parser')
        for script in soup(["script", "style", "nav", "footer"]):
            script.decompose()
        return soup.get_text(separator=' ', strip=True)
    except Exception: return None


# 📄 file readers (tuple yielding)
# ==========================================

# all readers yield (text_content, date_str, category_str)

def read_rows_raw_lines(file_bytes: bytes, encoding_choice: str = "auto") -> Iterable[Tuple[str, None, None]]:
    def _iter(enc):
        bio = io.BytesIO(file_bytes)
        with io.TextIOWrapper(bio, encoding=enc, errors="replace", newline=None) as wrapper:
            for line in wrapper: yield (line.rstrip("\r\n"), None, None)
    try:
        if encoding_choice == "latin-1": yield from _iter("latin-1")
        else: yield from _iter("utf-8")
    except UnicodeDecodeError:
        yield ("", None, None)

def read_rows_vtt(file_bytes: bytes, encoding_choice: str = "auto") -> Iterable[Tuple[str, None, None]]:
    # robust VTT reader that yields tuples
    def _iter_lines(enc):
        bio = io.BytesIO(file_bytes)
        with io.TextIOWrapper(bio, encoding=enc, errors="replace", newline=None) as wrapper:
            for line in wrapper: yield line.rstrip("\r\n")
    
    iterator = _iter_lines("utf-8") if encoding_choice != "latin-1" else _iter_lines("latin-1")
    
    for line in iterator:
        line = line.strip()
        if not line or line == "WEBVTT" or "-->" in line or line.isdigit(): continue
        if ":" in line:
            parts = line.split(":", 1)
            if len(parts) > 1 and len(parts[0]) < MAX_SPEAKER_NAME_LENGTH and " " in parts[0]:
                yield (parts[1].strip(), None, None)
                continue
        yield (line, None, None)

def read_rows_pdf(file_bytes: bytes) -> Iterable[Tuple[str, None, None]]:
    if pypdf is None: 
        st.error("pypdf missing")
        return
    bio = io.BytesIO(file_bytes)
    try:
        reader = pypdf.PdfReader(bio)
        for page in reader.pages:
            text = page.extract_text()
            if text: yield (text, None, None)
    except Exception: yield ("", None, None)

def read_rows_pptx(file_bytes: bytes) -> Iterable[Tuple[str, None, None]]:
    if pptx is None: return
    bio = io.BytesIO(file_bytes)
    try:
        prs = pptx.Presentation(bio)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                    if shape.text: yield (shape.text, None, None)
    except Exception:
        yield ("", None, None)

def read_rows_json(file_bytes: bytes, selected_key: str = None) -> Iterable[Tuple[str, None, None]]:
    bio = io.BytesIO(file_bytes)
    try:
        wrapper = io.TextIOWrapper(bio, encoding="utf-8", errors="replace")
        for line in wrapper:
            if not line.strip(): continue
            try:
                obj = json.loads(line)
                txt = ""
                if selected_key and isinstance(obj, dict): txt = str(obj.get(selected_key, ""))
                elif isinstance(obj, str): txt = obj
                else: txt = str(obj)
                yield (txt, None, None)
            except json.JSONDecodeError:
                pass
    except Exception:
        pass

def read_rows_csv_structured(
    file_bytes: bytes, 
    encoding_choice: str, 
    delimiter: str, 
    has_header: bool, 
    text_cols: List[str],
    date_col: Optional[str],
    cat_col: Optional[str],
    join_with: str
) -> Iterable[Tuple[str, Optional[str], Optional[str]]]:
    
    enc = "latin-1" if encoding_choice == "latin-1" else "utf-8"
    bio = io.BytesIO(file_bytes)
    
    with io.TextIOWrapper(bio, encoding=enc, errors="replace", newline="") as wrapper:
        rdr = csv.reader(wrapper, delimiter=delimiter)
        first = next(rdr, None)
        if first is None: return

        if has_header:
            header = make_unique_header(list(first))
            name_to_idx = {n: i for i, n in enumerate(header)}
            
            text_idxs = [name_to_idx[n] for n in text_cols if n in name_to_idx]
            date_idx = name_to_idx.get(date_col) if date_col else None
            cat_idx = name_to_idx.get(cat_col) if cat_col else None
        else:
            # if no header, user likely selected "col_0", "col_1" etc.
            name_to_idx = {f"col_{i}": i for i in range(len(first))}
            text_idxs = [name_to_idx[n] for n in text_cols if n in name_to_idx]
            date_idx = name_to_idx.get(date_col) if date_col else None
            cat_idx = name_to_idx.get(cat_col) if cat_col else None
            
            # yield 1st row data
            txt_parts = [first[i] if i < len(first) else "" for i in text_idxs]
            d_val = first[date_idx] if (date_idx is not None and date_idx < len(first)) else None
            c_val = first[cat_idx] if (cat_idx is not None and cat_idx < len(first)) else None
            yield (join_with.join(txt_parts), d_val, c_val)

        for row in rdr:
            txt_parts = [row[i] if i < len(row) else "" for i in text_idxs]
            d_val = row[date_idx] if (date_idx is not None and date_idx < len(row)) else None
            c_val = row[cat_idx] if (cat_idx is not None and cat_idx < len(row)) else None
            yield (join_with.join(txt_parts), d_val, c_val)

def iter_excel_structured(
    file_bytes: bytes, 
    sheet_name: str, 
    has_header: bool, 
    text_cols: List[str],
    date_col: Optional[str],
    cat_col: Optional[str],
    join_with: str
) -> Iterable[Tuple[str, Optional[str], Optional[str]]]:
    if openpyxl is None: return
    bio = io.BytesIO(file_bytes)
    wb = openpyxl.load_workbook(bio, read_only=True, data_only=True)
    ws = wb[sheet_name]
    rows_iter = ws.iter_rows(values_only=True)
    
    first = next(rows_iter, None)
    if first is None: 
        wb.close()
        return

    # header logic
    if has_header:
        header = make_unique_header(list(first))
        name_to_idx = {n: i for i, n in enumerate(header)}
        text_idxs = [name_to_idx[n] for n in text_cols if n in name_to_idx]
        date_idx = name_to_idx.get(date_col) if date_col else None
        cat_idx = name_to_idx.get(cat_col) if cat_col else None
    else:
        name_to_idx = {f"col_{i}": i for i in range(len(first))}
        text_idxs = [name_to_idx[n] for n in text_cols if n in name_to_idx]
        date_idx = name_to_idx.get(date_col) if date_col else None
        cat_idx = name_to_idx.get(cat_col) if cat_col else None
        
        # Yield first row
        txt_parts = [str(first[i]) if (first[i] is not None and i < len(first)) else "" for i in text_idxs]
        d_val = first[date_idx] if (date_idx is not None and date_idx < len(first)) else None
        c_val = first[cat_idx] if (cat_idx is not None and cat_idx < len(first)) else None
        yield (join_with.join(txt_parts), str(d_val) if d_val else None, str(c_val) if c_val else None)

    for row in rows_iter:
        txt_parts = [str(row[i]) if (row[i] is not None and i < len(row)) else "" for i in text_idxs]
        d_val = row[date_idx] if (date_idx is not None and date_idx < len(row)) else None
        c_val = row[cat_idx] if (cat_idx is not None and cat_idx < len(row)) else None
        yield (join_with.join(txt_parts), str(d_val) if d_val else None, str(c_val) if c_val else None)
    
    wb.close()

def detect_csv_headers(file_bytes: bytes, delimiter: str = ",") -> List[str]:
    try:
        bio = io.BytesIO(file_bytes)
        with io.TextIOWrapper(bio, encoding="utf-8", errors="replace", newline="") as wrapper:
            rdr = csv.reader(wrapper, delimiter=delimiter)
            row = next(rdr, None)
            return make_unique_header(row) if row else []
    except: return []

def detect_csv_num_cols(file_bytes: bytes, delimiter: str = ",") -> int:
    try:
        bio = io.BytesIO(file_bytes)
        with io.TextIOWrapper(bio, encoding="utf-8", errors="replace", newline="") as wrapper:
            rdr = csv.reader(wrapper, delimiter=delimiter)
            row = next(rdr, None)
            return len(row) if row else 0
    except: return 0

def get_excel_sheetnames(file_bytes: bytes) -> List[str]:
    if openpyxl is None: return []
    bio = io.BytesIO(file_bytes)
    wb = openpyxl.load_workbook(bio, read_only=True, data_only=True)
    sheets = list(wb.sheetnames)
    wb.close()
    return sheets

def get_excel_preview(file_bytes: bytes, sheet_name: str, has_header: bool, rows: int = 5) -> pd.DataFrame:
    if openpyxl is None: return pd.DataFrame()
    bio = io.BytesIO(file_bytes)
    try:
        df = pd.read_excel(bio, sheet_name=sheet_name, header=0 if has_header else None, nrows=rows, engine='openpyxl')
        if not has_header: df.columns = [f"col_{i}" for i in range(len(df.columns))]
        return df
    except:
        return pd.DataFrame()

def excel_estimate_rows(file_bytes: bytes, sheet_name: str, has_header: bool) -> int:
    if openpyxl is None: return 0
    bio = io.BytesIO(file_bytes)
    wb = openpyxl.load_workbook(bio, read_only=True, data_only=True)
    ws = wb[sheet_name]
    total = ws.max_row or 0
    wb.close()
    if has_header and total > 0: total -= 1
    return max(total, 0)


# ⚙️ processing logic
# ==========================================

def clean_date_str(raw: Any) -> Optional[str]:
    """Tries to extract a YYYY-MM-DD string from raw input."""
    if not raw: return None
    s = str(raw).strip()
    # somelightweight heuristics for ISO-like dates or standard formats
    # match YYYY-MM-DD
    match = re.search(r'\d{4}-\d{2}-\d{2}', s)
    if match: return match.group(0)
    # match MM/DD/YYYY
    match_us = re.search(r'(\d{1,2})/(\d{1,2})/(\d{4})', s)
    if match_us:
        # convert to ISO for sorting
        m, d, y = match_us.groups()
        return f"{y}-{int(m):02d}-{int(d):02d}"
    return None

def apply_text_cleaning(text: str, config: CleaningConfig) -> str:
    if not isinstance(text, str): return ""
    if config.remove_chat: text = CHAT_ARTIFACT_RE.sub(" ", text)
    if config.remove_html: text = HTML_TAG_RE.sub(" ", text)
    if config.unescape:
        try: text = html.unescape(text)
        except: pass
    if config.remove_urls: text = URL_EMAIL_RE.sub(" ", text)
    text = text.lower()
    if config.phrase_pattern: text = config.phrase_pattern.sub(" ", text)
    return text.strip()

def process_chunk_iter(
    rows_iter: Iterable[Tuple[str, Optional[str], Optional[str]]],
    clean_conf: CleaningConfig,
    proc_conf: ProcessingConfig,
    scanner: StreamScanner,
    lemmatizer: Optional[WordNetLemmatizer],
    progress_cb: Optional[Callable[[int], None]] = None
):
    _min_len = proc_conf.min_word_len
    _drop_int = proc_conf.drop_integers
    _trans = proc_conf.translate_map
    _stopwords = proc_conf.stopwords
    _lemma = proc_conf.use_lemmatization and (lemmatizer is not None)
    
    local_global_counts = Counter()
    local_global_bigrams = Counter() if proc_conf.compute_bigrams else Counter()
    
    batch_accum = Counter()
    batch_rows = 0
    row_count = 0
    
    # pre-caching lemmatizer methods for speed loop
    lemmatize = lemmatizer.lemmatize if _lemma else None

    for (raw_text, date_val, cat_val) in rows_iter:
        row_count += 1
        
        # entities (Before lowercase)
        if raw_text:
            entities = extract_entities_regex(raw_text, _stopwords)
            scanner.update_entities(entities)

        # cleaning
        text = apply_text_cleaning(raw_text, clean_conf)
        
        # tokenization & filter
        filtered_tokens_line: List[str] = []
        for t in text.split():
            t = t.translate(_trans)
            if not t: continue
            if _drop_int and t.isdigit(): continue
            if len(t) < _min_len: continue
            
            # lemmatize?
            if _lemma:
                # naive lemmatization (check if verb first, then noun) covers most bases without full part-of-speech tagging overhead
                t = lemmatize(t, pos='v')
                t = lemmatize(t, pos='n')
            
            if t in _stopwords: continue
            filtered_tokens_line.append(t)
        
        if filtered_tokens_line:
            # Stats Update
            line_counts = Counter(filtered_tokens_line)
            local_global_counts.update(filtered_tokens_line)
            
            if proc_conf.compute_bigrams and len(filtered_tokens_line) > 1:
                local_global_bigrams.update(pairwise(filtered_tokens_line))
            
            # metadata update
            clean_date = clean_date_str(date_val)
            clean_cat = str(cat_val).strip() if cat_val else None
            scanner.update_metadata_stats(clean_date, clean_cat, filtered_tokens_line)

            # topic modeling batching
            batch_accum.update(line_counts)
            batch_rows += 1
            if batch_rows >= scanner.DOC_BATCH_SIZE:
                scanner.add_topic_sample(batch_accum)
                batch_accum = Counter()
                batch_rows = 0

        if progress_cb and (row_count % 2000 == 0): progress_cb(row_count)

    # flushing last batch
    if batch_accum and batch_rows > 0:
        scanner.add_topic_sample(batch_accum)

    scanner.update_global_stats(local_global_counts, local_global_bigrams, row_count)
    if progress_cb: progress_cb(row_count)
    gc.collect()

def perform_refinery_job(file_obj, chunk_size, clean_conf: CleaningConfig):
    with tempfile.TemporaryDirectory() as temp_dir:
        original_name = os.path.splitext(file_obj.name)[0]
        status_container = st.status(f"⚙️ Refining {file_obj.name}...", expanded=True)
        part_num = 1
        created_files = []
        
        try:
            file_obj.seek(0)
            df_iterator = pd.read_csv(file_obj, chunksize=chunk_size, on_bad_lines='skip', dtype=str)
            
            for chunk in df_iterator:
                for col in chunk.columns:
                    chunk[col] = chunk[col].fillna("")
                    chunk[col] = chunk[col].apply(lambda x: apply_text_cleaning(x, clean_conf))
                
                new_filename = f"{original_name}_cleaned_part_{part_num}.csv"
                temp_path = os.path.join(temp_dir, new_filename)
                chunk.to_csv(temp_path, index=False)
                created_files.append(temp_path)
                status_container.write(f"✅ Processed chunk {part_num} ({len(chunk)} rows)")
                part_num += 1
            
            zip_buffer = io.BytesIO()
            with zipfile.ZipFile(zip_buffer, "w", zipfile.ZIP_DEFLATED) as zip_file:
                for file_path in created_files:
                    zip_file.write(file_path, arcname=os.path.basename(file_path))
            
            zip_buffer.seek(0)
            status_container.update(label="🎉 Refinery Job Complete!", state="complete", expanded=False)
            return zip_buffer
            
        except Exception as e:
            status_container.update(label="❌ Error", state="error")
            st.error(f"Refinery Error: {str(e)}")
            return None


# 📊 UI, analytics renderers
# ==========================================

def calculate_text_stats(counts: Counter, total_rows: int) -> Dict:
    total_tokens = sum(counts.values())
    unique_tokens = len(counts)
    avg_len = sum(len(word) * count for word, count in counts.items()) / total_tokens if total_tokens else 0
    return {
        "Total Rows": total_rows, "Total Tokens": total_tokens,
        "Unique Vocabulary": unique_tokens, "Avg Word Length": round(avg_len, 2),
        "Lexical Diversity": round(unique_tokens / total_tokens, 4) if total_tokens else 0
    }

def calculate_npmi(bigram_counts: Counter, unigram_counts: Counter, total_words: int, min_freq: int = 3) -> pd.DataFrame:
    results = []
    if not bigram_counts: return pd.DataFrame(columns=["Bigram", "Count", "NPMI"])
    epsilon = 1e-10 
    for (w1, w2), freq in bigram_counts.items():
        if freq < min_freq: continue
        count_w1 = unigram_counts.get(w1, 0)
        count_w2 = unigram_counts.get(w2, 0)
        if count_w1 == 0 or count_w2 == 0: continue
        prob_bigram = freq / total_words
        try: 
            pmi = math.log(prob_bigram / ((count_w1 / total_words) * (count_w2 / total_words)))
        except ValueError: continue
        
        log_prob_bigram = math.log(prob_bigram)
        if abs(log_prob_bigram) < epsilon: npmi = 1.0
        else: npmi = pmi / -log_prob_bigram
        results.append({"Bigram": f"{w1} {w2}", "Count": freq, "NPMI": round(npmi, 3)})
        
    df = pd.DataFrame(results)
    if df.empty: return pd.DataFrame(columns=["Bigram", "Count", "NPMI"])
    return df.sort_values("NPMI", ascending=False)

def calculate_tfidf(scanner: StreamScanner, top_n=50) -> pd.DataFrame:
    # IDF = log(total docs / doc freq)
    # TF (Global) = total count / total words (simplified for sketch)
    total_docs = len(scanner.topic_docs)
    if total_docs == 0: return pd.DataFrame()
    
    results = []
    # Analyze only terms that appear in at least 2 docs to avoid noise
    candidates = [t for t, c in scanner.doc_freqs.items() if c > 1]
    
    for term in candidates:
        tf = scanner.global_counts[term]
        df = scanner.doc_freqs[term]
        idf = math.log(total_docs / (1 + df))
        score = tf * idf
        results.append({"Term": term, "TF (Count)": tf, "DF (Docs)": df, "Keyphrase Score": round(score, 2)})
        
    df = pd.DataFrame(results)
    if df.empty: return df
    return df.sort_values("Keyphrase Score", ascending=False).head(top_n)

def perform_topic_modeling(synthetic_docs: List[Counter], n_topics: int, model_type: str) -> Optional[List[Dict]]:
    if not DictVectorizer or len(synthetic_docs) < 1: return None
    vectorizer = DictVectorizer(sparse=True)
    dtm = vectorizer.fit_transform(synthetic_docs)
    n_samples, n_features = dtm.shape
    if n_samples == 0 or n_features == 0: return None
    
    safe_n_topics = min(n_topics, min(n_samples, n_features)) if model_type == "NMF" else min(n_topics, n_samples)
    if safe_n_topics < 1: return None

    model = None
    try:
        if model_type == "LDA": model = LatentDirichletAllocation(n_components=safe_n_topics, random_state=42, max_iter=10)
        elif model_type == "NMF": model = NMF(n_components=safe_n_topics, random_state=42, init='nndsvd')
        model.fit(dtm)
    except ValueError: return None
    
    feature_names = vectorizer.get_feature_names_out()
    topics = []
    for topic_idx, topic in enumerate(model.components_):
        top_indices = topic.argsort()[:-7:-1]
        top_words = [feature_names[i] for i in top_indices]
        strength = sum(topic[i] for i in top_indices)
        topics.append({"id": topic_idx + 1, "words": top_words, "strength": strength})
    return topics

def perform_bayesian_sentiment_analysis(counts: Counter, sentiments: Dict[str, float], pos_thresh: float, neg_thresh: float) -> Optional[Dict]:
    if not beta_dist: return None
    pos_count = sum(counts[w] for w, s in sentiments.items() if s >= pos_thresh)
    neg_count = sum(counts[w] for w, s in sentiments.items() if s <= neg_thresh)
    total_informative = pos_count + neg_count
    if total_informative < 1: return None

    alpha_post = 1 + pos_count
    beta_post = 1 + neg_count
    mean_prob = alpha_post / (alpha_post + beta_post)
    lower_ci, upper_ci = beta_dist.ppf([0.025, 0.975], alpha_post, beta_post)
    x = np.linspace(0, 1, 300)
    y = beta_dist.pdf(x, alpha_post, beta_post)
    return {
        "pos_count": pos_count, "neg_count": neg_count, "total": total_informative,
        "mean_prob": mean_prob, "ci_low": lower_ci, "ci_high": upper_ci,
        "x_axis": x, "pdf_y": y
    }

def render_workflow_guide():
    with st.expander("📘 Comprehensive App Guide: How to use this Tool", expanded=False):
        st.markdown("""
        ### 🌟 What is this?
        **Signal Foundry is an Unstructured Data Intelligence Engine (or: an "MRI for your documents").** 
        It takes raw, "dirty" text and extracts mathematical structure **(Engineering View)** to reveal hidden risks, patterns, and causal links **(Value View)** without requiring you to write code.

        ---

        ### 🚀 1. Choose Your Workflow

        #### A. The "Direct Scan" (Standard)
        *   **Best for:** PDFs, PowerPoints, Transcripts, or CSVs/JSON files.
        *   **How:** 
            1. Upload your files in the sidebar.
            2. **Crucial Step:** A "Scan Configuration" box will appear. 
            3. **For Trends:** Select your **Date Column** (e.g., for logs or news).
            4. Click **"⚡ Start Scan"**.
        *   **Result:** A complete statistical "Sketch" of the document's DNA.

        #### B. The "Enterprise" Workflow (Secure / Offline)
        *   **Best for:** Massive datasets (>10M rows) or sensitive PII that cannot leave a secure server.
        *   **How:** Use the offline `harvester.py` script to generate a `.json` Sketch file (containing only math/counts, no raw text). Upload that JSON here to visualize it safely.

        ---

        ### 🧠 2. Interpret the Analytics (The Foundry Process)

        #### 📈 Technical Signal / Keyphrases (TF-IDF)
        *   **The Tech:** Uses *Term Frequency-Inverse Document Frequency* to filter out generic high-frequency words (like "report", "therefore").
        *   **The Value:** Isolates the **Unique Signal**. It answers: *"What is the specific jargon or technical topic unique to this file?"*

        #### 🔬 Sticky Concepts / Phrase Significance (NPMI)
        *   **The Tech:** Uses *Normalized Pointwise Mutual Information* to measure the mathematical probability of two words occurring together vs. random chance.
        *   **The Value:** Reveals **Specific Risks**. It distinguishes generic "Data" from specific "Biometric Data" or "Dual Use." It finds the "terms of art."

        #### 👥 Stakeholder Map / Entities (NER Lite)
        *   **The Tech:** Heuristic extraction of Capitalized Phrases.
        *   **The Value:** Identifies **Who and What**. Quickly maps the Actors (Organizations, People) and Regulations cited in the text.

        #### 🕸️ Context & Causality / Network Graph
        *   **The Tech:** A force-directed graph connecting nodes based on physical proximity in the text.
        *   **The Value:** Reveals **Hidden Associations**. If "Stimulation" links to "Therapy," the context is medical. If it links to "Interrogation," the context is weaponization.

        #### 🔍 Thematic Buckets / Topic Modeling
        *   **The Tech:** Uses Matrix Factorization (NMF) or Probabilistic Models (LDA) to group words into vectors.
        *   **The Value:** automated **Categorization**.
            *   **LDA:** Best for long reports (mixed topics).
            *   **NMF:** Best for short chats/tickets (distinct topics).

        #### ⚖️ Risk Confidence / Bayesian Sentiment
        *   **The Tech:** Calculates a *Credible Interval* using Beta Distributions rather than a single raw score.
        *   **The Value:** **Decision Safety**. Instead of saying "Sentiment is 60%," it says "We are 95% confident the true sentiment lies between 55% and 65%." This prevents knee-jerk reactions to statistically insignificant data.
        """)

def render_lit_case_study():
    # We use Unicode "Math Sans" characters to simulate bold/italics in the title
    # Italic 'another': 𝘢𝘯𝘰𝘵𝘩𝘦𝘳
    # Bold 'specific': 𝘀𝗽𝗲𝗰𝗶𝗳𝗶𝗰
    title = "🔦 Spotlight: Digital Humanities & Literary Forensics (𝘢𝘯𝘰𝘵𝘩𝘦𝘳 𝘀𝗽𝗲𝗰𝗶𝗳𝗶𝗰 Case Study)"
    
    with st.expander(title, expanded=False):
        st.markdown("""
        ### The Scenario
        **The Artifact:** The full text of Ovid's **<a href="https://www.gutenberg.org/files/21765/21765-h/21765-h.htm" target="_blank">"Metamorphoses"</a>** (via Project Gutenberg URL).
        **The User:** A Digital Humanities Researcher or Student.
        **The Goal:** To rapidly map the "Pantheon" of characters and distinguish the original narrative from the translator's artifacts.

        ---

        ### 1. The "Pantheon Map" (Entities Tab)
        *   **The Question:** "Who are the dominant power players in this 15-book epic?"
        *   **The Signal:** Capitalized Name Extraction.
        *   **The Result:** The engine immediately surfaces **"Jupiter," "Apollo," "Ceres,"** and **"Minerva"** as the top nodes.
        *   **The Value / Insight:** Without reading a single line, you have a hierarchical map of the Roman deities driving the plot.

        ### 2. The "Translator's Fingerprint" (NPMI & Bigrams)
        *   **The Question:** "Is this pure text, or is there structural noise?"
        *   **The Signal:** Sticky Concepts (Bigrams).
        *   **The Result:** The engine identifies **"Clarke translates"** and **"-ver Clarke"** as top phrases.
        *   **The Value / Insight:** **Forensic Separation.** The engine detected that *John Clarke* (the translator) is statistically inseparable from the text. It highlights "Data Hygiene" issues—showing you exactly what "boilerplate" needs to be cleaned (e.g., "Project Gutenberg" headers) before deep analysis.

        ### 3. The "Narrative Arcs" (Topic Modeling)
        *   **The Question:** "What are the distinct recurring themes?"
        *   **The Signal:** NMF/LDA Mathematical Bucketing.
        *   **The Result:**
            *   **Topic A:** [Daughter, Jupiter, Cadmus, Wife] -> *The Genealogy & Origin Myths.*
            *   **Topic B:** [Thou, Thee, Thus, Said] -> *The Dialogue & Poetic Structure.*
        *   **The Value / Insight:** The engine successfully separates the *Style* (Archaic English) from the *Substance* (Mythological Events).

        ### 4. The "Semantic Network" (Graph Tab)
        *   **The Question:** "How do the main characters interact?"
        *   **The Signal:** Proximity-based linking.
        *   **The Result:** "jupiter" is the central "hub" node, with spokes connecting to various "nymphs" and "daughters."
        *   **The Value / Insight:** Visualizes the centralized power structure of the mythology, confirming Jupiter as the primary driver of the transformations.
        """, unsafe_allow_html=True)

def render_auto_insights(scanner, proc_conf):
    # Only run if we have data
    if not scanner.global_counts: return

    # --- 1. PREPARE DATA ---
    # Entities
    top_ents = scanner.entity_counts.most_common(3)
    ent_str = ", ".join([f"**{e[0]}**" for e in top_ents]) if top_ents else "(No entities detected)"
    
    # Sticky Concepts (NPMI)
    df_npmi = calculate_npmi(scanner.global_bigrams, scanner.global_counts, scanner.total_rows_processed)
    top_npmi = df_npmi.head(3)["Bigram"].tolist() if not df_npmi.empty else []
    npmi_str = ", ".join([f"**{b}**" for b in top_npmi]) if top_npmi else "(No strong phrases found)"
    
    # Tech Signal (TF-IDF)
    df_tfidf = calculate_tfidf(scanner, 20)
    top_idf = df_tfidf.head(3)["Term"].tolist() if not df_tfidf.empty else []
    idf_str = ", ".join([f"**{t}**" for t in top_idf]) if top_idf else "(Not enough documents for TF-IDF)"

    # -render reporting
    with st.expander("⚡ High-Level Signal Report (Auto-Generated, ymmv)", expanded=True):
        st.markdown(f"""
        ### 1. The "Stakeholder Map" (Entities)
        *   **The Question:** "Who are the dominant actors or organizations?"
        *   **The Signal:** Capitalized Name Extraction.
        *   **The Result:** {ent_str}
        *   **The Insight:** These nodes appear most frequently, suggesting they are the primary drivers of the narrative or the key subjects of the file.

        ### 2. The "Sticky Concepts" (Phrase Significance)
        *   **The Question:** "What is the specific 'Term of Art' or jargon here?"
        *   **The Signal:** NPMI (Normalized Pointwise Mutual Information).
        *   **The Result:** {npmi_str}
        *   **The Insight:** These words appear together mathematically more often than random chance, indicating they represent specific concepts (e.g. "Credit Card" vs "Red Card") rather than generic language.

        ### 3. The "Technical Signal" (Keyphrases)
        *   **The Question:** "What makes this specific document unique?"
        *   **The Signal:** TF-IDF (Inverse Document Frequency).
        *   **The Result:** {idf_str}
        *   **The Insight:** While words like "the" or "report" might be frequent, *these* specific words are statistically unique to this dataset, representing its core technical signature.
        """)

def render_neurotech_case_study():
    with st.expander("🔦 Spotlight: Analyzing Mi|itary Neurotechno|ogy (a very *specific* Case Study)", expanded=False):
        st.markdown("""
        ### The Scenario
        **The Artifact:** A dense, 50-page UNIDIR report titled <a href="https://unidir.org/wp-content/uploads/2025/11/UNIDIR_Neurotechnology_Military-Domain_A-Primer.pdf" target="_blank"><b>"Neurotechnology in the Military Domain"</b></a>.
        **The User:** A Defense Ana|yst with 5 minutes to extract actionable inte||igence.
        **The Goal:** Move beyond "what is this paper about?" to "what are the threats and opportunities?"

        ---

        ### 1. The "Sticky Concepts" (NPMI Tab)
        *   **The Question:** "What specific types of risks are discussed?"
        *   **The Signal:** The engine finds words that mathematically *stick together* more than random chance.
        *   **The Result:** It surfaces **"Dua| Use"** and **"Cognitive Liberty."**
        *   **The Insight:** The strategic risk isn't just new weap0ns; it is civi|ian medica| techno|ogy being repurposed for mi|itary app|ications (Dua| Use), necessitating a legal/ethical framework (Liberty).

        ### 2. The "Technical Signal" (Keyphrases Tab)
        *   **The Question:** "Do I need to worry about brain implants yet?"
        *   **The Signal:** TF-IDF filters out generic words to find unique technical terms.
        *   **The Result:** High scores for **"Non-invasive,"** **"Transcranial,"** and **"Wearable."**
        *   **The Insight:** The immediate operationa| reality is external headsets/helmets, not surgical implants.

        ### 3. The "Semantic Network" (Graph Tab)
        *   **The Question:** "How is the techno|ogy being applied?"
        *   **The Signal:** The Graph links words based on proximity in the text.
        *   **The Result:** 
            *   Cluster A links **"Stimulation"** to **"Performance"** (Enhancement/Super-Soldiers).
            *   Cluster B links **"Stimulation"** to **"Interrogation"** (Weap0nization/T0rture).
        *   **The Insight:** The paper treats "Enhancement" and "Weaponization" as distinct operational clusters.

        ### 4. The "Stakeholder Map" (Entities Tab)
        *   **The Question:** "Who is involved?"
        *   **The Signal:** Capitalized Name Extraction.
        *   **The Result:** **"D@RPA," "Neura|ink," "Geneva Convention," "Human Rights Council."**
        *   **The Insight:** Identifies the funding sources (DARP@) vs. the regulatory blockers (Geneva).
        """, unsafe_allow_html=True)


def render_use_cases():
    with st.expander("📖 (some more general) Use-cases", expanded=False):
        st.markdown("""
        ### (Some) use-cases for this unstructured data intelligence engine; you'll think of more...
        
        #### 📈 Temporal & Trends
        *   **Crisis Timeline:** Map high-severity words (e.g., "leak", "fail") over time to pinpoint incident starts.
        *   **Campaign Tracking:** Watch how slogan adoption grows or fades week-over-week.
        
        #### 🏢 Corporate & Strategic
        *   **Stakeholder Mapping (NER):** Instantly see *Who* and *What* are driving the conversation.
        *   **Customer Feedback:** "Diffing" analysis to see what Engineering words are distinct from Sales words.
        *   **M&A Due Diligence:** Scanning Data Rooms for liability terms without reading files.

        #### 🔬 Research & Forensics
        *   **Signal vs Noise:** TF-IDF extraction to ignore frequent words and find unique cluster identifiers.
        *   **Literary Forensics:** Vocabulary diversity and phrase patterns.
        *   **Unknown Unknowns:** Surface recurring challenges via bigram maps.

        #### 🛡️ Security & Privacy
        *   **The "Privacy Proxy":** Refining data locally before sending sanitized stats to LLMs.
        *   **Insider Threat:** Pattern matching on communication logs.
        
        ---
        #### 🎓 Education: LMS Discussion Forums
        *Insights into what a group of students are discussing:*
        *   Identifying the most common topics and themes students are talking about.
        *   Surfacing frequently asked questions and recurring challenges.
        *   Detecting sentiment trends (e.g., frustration, excitement, confusion) across the class.
        *   Visualizing connections between concepts or issues using network graphs.
        *   Highlighting emerging issues or "unknown unknowns" (e.g., a misunderstood assignment).
        *   Comparing discussion dynamics before and after key events (e.g., exams).
        *   Summarizing participation patterns (who is most/least active).
        *   Providing instructors with actionable summaries for targeted intervention.
        """)


def render_analyst_help():
    with st.expander("🎓 Analyst's Guide", expanded=False):
        st.markdown("""
        **Symptom: Graph vs. Topics Disagree**
        *   **Fix:** Check for 'bridge' words. Lower 'Rows per Document' to 1.
        
        **Symptom: Giant Blob Graph**
        *   **Fix:** Increase 'Min Link Frequency'.
        """)

# visualization helpers
@st.cache_data(show_spinner="Analyzing term sentiment...")
def get_sentiments(_analyzer, terms: Tuple[str, ...]) -> Dict[str, float]:
    if not _analyzer or not terms: return {}
    return {term: _analyzer.polarity_scores(term)['compound'] for term in terms}

def create_sentiment_color_func(sentiments: Dict[str, float], pos_color, neg_color, neu_color, pos_thresh, neg_thresh):
    def color_func(word, font_size, position, orientation, random_state=None, **kwargs):
        score = sentiments.get(word, 0.0)
        if score >= pos_thresh: return pos_color
        elif score <= neg_thresh: return neg_color
        else: return neu_color
    return color_func

def get_sentiment_category(score: float, pos_threshold: float, neg_threshold: float) -> str:
    if score >= pos_threshold: return "Positive"
    if score <= neg_threshold: return "Negative"
    return "Neutral"

def build_wordcloud_figure_from_counts(counts, max_words, width, height, bg_color, colormap, font_path, random_state, color_func):
    limited = dict(counts.most_common(max_words))
    if not limited: return plt.figure(), None
    wc = WordCloud(width=width, height=height, background_color=bg_color, colormap=colormap, font_path=font_path, random_state=random_state, color_func=color_func, collocations=False, normalize_plurals=False).generate_from_frequencies(limited)
    fig, ax = plt.subplots(figsize=(max(6, width/100), max(3, height/100)), dpi=100)
    ax.imshow(wc, interpolation="bilinear"); ax.axis("off"); plt.tight_layout()
    return fig, wc

def fig_to_png_bytes(fig):
    buf = BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight", pad_inches=0.1)
    buf.seek(0)
    return buf

# 🤖 AI logic
def call_llm_and_track_cost(system_prompt: str, user_prompt: str, config: dict):
    try:
        client = openai.OpenAI(api_key=config['api_key'], base_url=config['base_url'])
        response = client.chat.completions.create(
            model=config['model_name'],
            messages=[
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ]
        )
        in_tok = 0
        out_tok = 0
        if hasattr(response, 'usage') and response.usage:
            in_tok = response.usage.prompt_tokens
            out_tok = response.usage.completion_tokens
        
        cost = (in_tok * config['price_in'] / 1_000_000) + (out_tok * config['price_out'] / 1_000_000)
        st.session_state['total_tokens'] += (in_tok + out_tok)
        st.session_state['total_cost'] += cost
        return response.choices[0].message.content
    except Exception as e:
        return f"AI Error: {str(e)}"


# 🚀 main app ui
# ==========================================

st.set_page_config(page_title="Signal Foundry", layout="wide")
st.toast("app loaded successfully", icon="🚀") # cache buster
st.title("🧠 Signal Foundry: Unstructured Data Analytics")
st.markdown("### *(or: data geiger counter~)*")

render_workflow_guide() # calling restored full guide
render_use_cases()
render_neurotech_case_study()
render_lit_case_study()
analyzer, lemmatizer = setup_nlp_resources()

# --- sidebar
with st.sidebar:
    st.header("📂 Data Input")
    uploaded_files = st.file_uploader("Upload Files", type=["csv", "xlsx", "vtt", "txt", "json", "pdf", "pptx"], accept_multiple_files=True)
    
    # --modfified logic
    
    # checking if there's currently data
    has_data = st.session_state['sketch'].total_rows_processed > 0
    
    # showing the checkbox
    clear_on_scan = st.checkbox("Clear previous data", value=False)
    
    # 'banner' logic: only showing if there's ambiguity
    if has_data and not clear_on_scan:
        st.info("⚠️ **Additive Mode Active:** New scans will be ADDED to current results. Check the box above to Start Fresh.", icon="ℹ️")
    elif has_data and clear_on_scan:
        st.caption("✅ Next scan will overwrite current data.")
        
    if st.button("🗑️ Reset All"): reset_sketch(); st.rerun()
    
    # --end
    
    st.divider()
    with st.expander("🌐 Web/Text Import"):
        sketch_upload = st.file_uploader("Import Sketch (.json)", type=["json"])
        if sketch_upload:
            file_hash = hash(sketch_upload.getvalue())
            if st.session_state.get('last_sketch_hash') != file_hash:
                if st.session_state['sketch'].load_from_json(sketch_upload.getvalue().decode('utf-8')):
                    st.session_state['last_sketch_hash'] = file_hash
                    st.success("Sketch Loaded Successfully!")
                else:
                    st.error("Invalid Sketch File")

        url_input = st.text_area("URLs (1 per line; doesn't always work with more modern sites)")
        manual_input = st.text_area("Manual Text")
        
    st.divider()
    st.header("🔐 AI Setup")
    if st.session_state['authenticated']:
        st.success("Unlocked")
        with st.expander("🤖 Provider Settings", expanded=True):
            ai_provider = st.radio("Provider", ["xAI (Grok)", "OpenAI (GPT-4o)"])
            if "OpenAI" in ai_provider:
                api_key_name = "openai_api_key"
                base_url = None 
                model_name = st.selectbox("Model", ["gpt-4o", "gpt-4o-mini"])
                price_in, price_out = (0.15, 0.60) if "mini" in model_name else (2.50, 10.00)
            else:
                api_key_name = "xai_api_key"
                base_url = "https://api.x.ai/v1"
                model_name = "grok-4-0709" 
                price_in, price_out = 3.00, 15.00
            
            api_key = st.secrets.get(api_key_name)
            if not api_key: api_key = st.text_input(f"Enter {api_key_name}", type="password")
            
            ai_config = {'api_key': api_key, 'base_url': base_url, 'model_name': model_name, 'price_in': price_in, 'price_out': price_out}

        with st.expander("💰 Cost Estimator", expanded=False):
            c1, c2 = st.columns(2)
            c1.markdown(f"**Tokens:**\n{st.session_state['total_tokens']:,}")
            c2.markdown(f"**Cost:**\n`${st.session_state['total_cost']:.5f}`")
            if st.button("Reset Cost"):
                st.session_state['total_cost'] = 0.0
                st.session_state['total_tokens'] = 0
                st.rerun()
        if st.button("Logout"): logout(); st.rerun()
    else:
        st.text_input("Password", type="password", key="password_input", on_change=perform_login)
        if st.session_state['auth_error']: st.error("Incorrect password")

    st.divider()
    st.header("⚙️ Configuration")
    
    st.markdown("**Cleaning**")
    clean_conf = CleaningConfig(
        remove_chat=st.checkbox("Remove Chat Artifacts", True, help="Strips metadata like timestamps, usernames (e.g., <@U1234>), and system messages from logs/transcripts to focus purely on the conversation content."),
        remove_html=st.checkbox("Remove HTML", True),
        remove_urls=st.checkbox("Remove URLs", True),
        unescape=st.checkbox("Unescape HTML", True, help="Converts coded entities (e.g., &amp ; amp;, &amp ; quot;) back into readable symbols (&, \").")
    )
    
    st.markdown("**Processing**")
    use_lemma = st.checkbox("Use Lemmatization", False, help="Merges 'running' -> 'run'. Slower but cleaner.")
    if use_lemma and lemmatizer is None: st.warning("NLTK Lemmatizer not found.")
    
    proc_conf = ProcessingConfig(
        min_word_len=st.slider("Min Word Len", 1, 10, 2),
        drop_integers=st.checkbox("Drop Integers", True),
        compute_bigrams=st.checkbox("Bigrams", True),
        use_lemmatization=use_lemma,
        translate_map=build_punct_translation(st.checkbox("Keep Hyphens"), st.checkbox("Keep Apostrophes"))
    )
    
    # stopwords
    user_sw = st.text_area("Stopwords (comma-separated)", "firstname.lastname, jane doe, okay, ok, really")
    phrases, singles = parse_user_stopwords(user_sw)
    clean_conf.phrase_pattern = build_phrase_pattern(phrases)
    stopwords = set(STOPWORDS).union(singles)
    if st.checkbox("Remove Prepositions", True): stopwords.update(default_prepositions())
    proc_conf.stopwords = stopwords
    
    st.markdown("### 🎨 Appearance")
    bg_color = st.color_picker("Background Color", "#000000")
    colormap = st.selectbox("Colormap", ["viridis", "plasma", "inferno", "magma", "cividis", "tab10", "Blues", "Reds", "Greys"], 0)
    top_n = st.number_input("Top Terms to Display", min_value=5, max_value=1000, value=20)
    max_words = st.slider("Max Words (Cloud)", 50, 3000, 1000, 50)
    
    # font selection
    font_map, font_names = list_system_fonts(), list(list_system_fonts().keys())
    default_font_idx = 0
    desired_font = "DejaVu Sans Mono"
    if desired_font in font_names:
        default_font_idx = font_names.index(desired_font)
    combined_font_name = st.selectbox("Font", font_names or ["(default)"], 0)
    combined_font_path = font_map.get(combined_font_name) if font_names else None

    st.markdown("### 🔬 Sentiment")
    enable_sentiment = st.checkbox("Enable Sentiment", False)
    if enable_sentiment and analyzer is None:
        st.error("NLTK not found.")
        enable_sentiment = False
    
    pos_threshold, neg_threshold, pos_color, neu_color, neg_color = 0.05, -0.05, '#2ca02c', '#808080', '#d62728'
    if enable_sentiment:
        c1, c2 = st.columns(2)
        with c1: pos_threshold = st.slider("pos threshold", 0.0, 1.0, 0.05, 0.01)
        with c2: neg_threshold = st.slider("neg threshold", -1.0, 0.0, -0.05, 0.01)
        c1, c2, c3 = st.columns(3)
        with c1: pos_color = st.color_picker("pos color", value=pos_color)
        with c2: neu_color = st.color_picker("neu color", value=neu_color)
        with c3: neg_color = st.color_picker("neg color", value=neg_color)

    # Updated with help text
    doc_granularity = st.select_slider(
        "Rows per Doc", 
        options=[1, 5, 10, 100, 500], 
        value=5,
        help="How many lines of text form one 'Document'. Use 1-5 for Chat Logs/Tweets. Use 100+ for Books/Reports."
    )
    st.session_state['sketch'].set_batch_size(doc_granularity)
    
    current_settings_hash = f"{doc_granularity}_{proc_conf.min_word_len}"
    if 'last_settings_hash' not in st.session_state: 
        st.session_state['last_settings_hash'] = current_settings_hash

    if st.session_state['last_settings_hash'] != current_settings_hash:
        if st.session_state['sketch'].total_rows_processed > 0:
            reset_sketch()
            st.warning("⚙️ Settings changed. Data reset for consistency. Please Scan again.")
        st.session_state['last_settings_hash'] = current_settings_hash

    # Updated with help text
    topic_model_type = st.selectbox(
        "Topic Model", 
        ["LDA", "NMF"],
        help="LDA: Best for long text with mixed topics (Probabilistic). NMF: Best for short text with distinct topics (Linear Algebra)."
    )
    
    # Updated with help text
    n_topics = st.slider(
        "Topics", 
        2, 10, 
        4,
        help="The number of distinct themes the AI should attempt to find."
    )

with st.expander("🛠️ Data Refinery"):
    ref_file = st.file_uploader("CSV to Refine", type=['csv'])
    if ref_file and st.button("🚀 Run Refinery"):
        zip_data = perform_refinery_job(ref_file, 50000, clean_conf)
        if zip_data: st.download_button("Download ZIP", zip_data, "refined.zip", "application/zip")

# --scanning phase
all_inputs = list(uploaded_files) if uploaded_files else []
if url_input:
    for u in url_input.split('\n'):
        if u.strip(): 
            txt = fetch_url_content(u.strip())
            if txt: 
                all_inputs.append(VirtualFile(f"url_{hash(u)}.txt", txt))
                time.sleep(URL_SCRAPE_RATE_LIMIT_SECONDS) # RATE LIMITING
if manual_input: all_inputs.append(VirtualFile("manual.txt", manual_input))

if all_inputs:
    st.subheader("🚀 Scanning Phase")
    
    for idx, f in enumerate(all_inputs):
        try:
            # resource limit check
            if f.getbuffer().nbytes > MAX_FILE_SIZE_MB * 1024 * 1024:
                st.error(f"❌ File **{f.name}** exceeds {MAX_FILE_SIZE_MB}MB limit.")
                continue

            file_bytes, fname, lower = f.getvalue(), f.name, f.name.lower()
            is_csv = lower.endswith(".csv")
            is_xlsx = lower.endswith((".xlsx", ".xlsm"))
            is_json = lower.endswith(".json")
            is_vtt = lower.endswith(".vtt")
            is_pdf = lower.endswith(".pdf")
            is_pptx = lower.endswith(".pptx")
            
            # Default Scan Settings
            scan_settings = {
                "date_col": None,
                "cat_col": None,
                "text_cols": [],
                "has_header": False,
                "sheet_name": None,
                "json_key": None
            }
            
            with st.expander(f"🧩 Config: {fname}", expanded=True):
                if is_csv:
                    headers = detect_csv_headers(file_bytes)
                    if headers:
                        scan_settings["has_header"] = True
                        st.info(f"Detected {len(headers)} columns.")
                        scan_settings["text_cols"] = st.multiselect("Text Columns", headers, default=[headers[0]], key=f"txt_{idx}")
                        scan_settings["date_col"] = st.selectbox("Date Column (Optional)", ["(None)"] + headers, key=f"date_{idx}")
                        scan_settings["cat_col"] = st.selectbox("Category Column (Optional)", ["(None)"] + headers, key=f"cat_{idx}")
                        
                        if scan_settings["date_col"] == "(None)": scan_settings["date_col"] = None
                        if scan_settings["cat_col"] == "(None)": scan_settings["cat_col"] = None
                    else:
                        st.warning("No headers detected. Scanning as raw text.")
                elif is_xlsx:
                    sheets = get_excel_sheetnames(file_bytes)
                    scan_settings["sheet_name"] = st.selectbox("Sheet", sheets, key=f"sheet_{idx}")
                    if scan_settings["sheet_name"]:
                         scan_settings["has_header"] = st.checkbox("Has Header Row", True, key=f"xls_head_{idx}")
                elif is_json:
                    scan_settings["json_key"] = st.text_input("JSON Key (Optional)", "", key=f"json_{idx}")

            if st.button(f"⚡ Start Scan: {fname}", key=f"btn_{idx}"):
                if clear_on_scan: reset_sketch()
                bar = st.progress(0)
                status = st.empty()
                
                # select iterator
                rows_iter = iter([])
                approx = estimate_row_count_from_bytes(file_bytes)

                if is_csv and scan_settings["has_header"] and scan_settings["text_cols"]:
                    rows_iter = read_rows_csv_structured(
                        file_bytes, "auto", ",", True, 
                        scan_settings["text_cols"], scan_settings["date_col"], scan_settings["cat_col"], " "
                    )
                elif is_xlsx and scan_settings["sheet_name"]:
                    rows_iter = iter_excel_structured(
                        file_bytes, scan_settings["sheet_name"], scan_settings["has_header"], 
                        ["col_0"], None, None, " " 
                    )
                elif is_pdf:
                    rows_iter = read_rows_pdf(file_bytes)
                elif is_pptx:
                    rows_iter = read_rows_pptx(file_bytes)
                elif is_vtt:
                    rows_iter = read_rows_vtt(file_bytes)
                elif is_json:
                    rows_iter = read_rows_json(file_bytes, scan_settings["json_key"])
                else:
                    # fallback raw line reader
                    rows_iter = read_rows_raw_lines(file_bytes)
                
                # run
                process_chunk_iter(rows_iter, clean_conf, proc_conf, st.session_state['sketch'], lemmatizer, lambda n: status.text(f"Rows: {n:,}"))
                bar.progress(100)
                status.success("Done!")
                if not clear_on_scan: st.rerun()

        except Exception as e:
            st.error(f"Error: {e}")

# --- analysis phase

scanner = st.session_state['sketch']

# now: dynamic filtering
# filter the view based on the CURRENT slider setting. 
# ensures that if user changes the slider from 2 to 7, the 
# wordcloud updates immediately without needing to re-scan the files
combined_counts = Counter({
    k: v for k, v in scanner.global_counts.items() 
    if len(str(k)) >= proc_conf.min_word_len
    and k not in proc_conf.stopwords # <--- New Logic
})

if combined_counts:
    st.divider()
    st.header("📊 Analysis Dashboard")
    
    # calculate stats upfront
    text_stats = calculate_text_stats(combined_counts, scanner.total_rows_processed)
    render_auto_insights(scanner, proc_conf)
    # main tabs
    tab_main, tab_trend, tab_ent, tab_key = st.tabs(["☁️ Word Cloud & Stats", "📈 Trends", "👥 Entities", "🔑 Keyphrases"])
    
    with tab_main:
        if enable_sentiment:
            top_keys = [k for k,v in combined_counts.most_common(1000)]
            term_sentiments = get_sentiments(analyzer, tuple(top_keys))
            if proc_conf.compute_bigrams:
                 top_bg_keys = [" ".join(k) for k,v in scanner.global_bigrams.most_common(2000)]
                 term_sentiments.update(get_sentiments(analyzer, tuple(top_bg_keys)))
            c_color_func = create_sentiment_color_func(term_sentiments, pos_color, neg_color, neu_color, pos_threshold, neg_threshold)
            fig, _ = build_wordcloud_figure_from_counts(combined_counts, max_words, 800, 400, bg_color, colormap, combined_font_path, 42, c_color_func)
        else:
            term_sentiments = {}
            fig, _ = build_wordcloud_figure_from_counts(combined_counts, max_words, 800, 400, bg_color, colormap, combined_font_path, 42, None)
            
        st.pyplot(fig, use_container_width=True)
        st.download_button("📥 download combined png", fig_to_png_bytes(fig), "combined_wc.png", "image/png")
        
        c1, c2, c3, c4 = st.columns(4)
        c1.metric("Total Tokens", f"{text_stats['Total Tokens']:,}")
        c2.metric("Unique Vocab", f"{text_stats['Unique Vocabulary']:,}")
        c3.metric("Docs/Rows", f"{text_stats['Total Rows']:,}")
        c4.metric("Lexical Diversity", f"{text_stats['Lexical Diversity']}")

    with tab_trend:
        if scanner.temporal_counts:
            st.markdown("#### Word Volume Over Time")
            trend_data = []
            for d_str, counts in scanner.temporal_counts.items():
                trend_data.append({"Date": d_str, "Volume": sum(counts.values())})
            
            df_trend = pd.DataFrame(trend_data).sort_values("Date")
            st.line_chart(df_trend.set_index("Date"))
            
            st.markdown("#### Specific Term Trends")
            terms_to_plot = st.multiselect("Select terms to plot", [t for t, c in combined_counts.most_common(50)])
            if terms_to_plot:
                term_trend_data = []
                for d_str, counts in scanner.temporal_counts.items():
                    row = {"Date": d_str}
                    for t in terms_to_plot: row[t] = counts[t]
                    term_trend_data.append(row)
                df_term_trend = pd.DataFrame(term_trend_data).sort_values("Date").set_index("Date")
                st.line_chart(df_term_trend)
        else:
            st.info("No Date column was selected during scan (or no valid dates found).")

    with tab_ent:
        st.markdown("#### Top Entities (NER Lite)")
        if scanner.entity_counts:
            ent_df = pd.DataFrame(scanner.entity_counts.most_common(50), columns=["Entity", "Count"])
            st.dataframe(ent_df, use_container_width=True)
            
            # simple entity cloud (safety wrapped)
            try:
                fig_e, _ = build_wordcloud_figure_from_counts(scanner.entity_counts, 100, 800, 400, "#111111", "Pastel1", combined_font_path, 42, None)
                st.pyplot(fig_e)
            except Exception as e:
                st.warning(f"Could not generate Entity Cloud: {e}")
        else:
            st.info("No capitalized entities detected.")

    with tab_key:
        st.markdown("#### TF-IDF Keyphrases")
        st.caption("These words are 'Unique' to specific documents, filtered out generic high-frequency noise.")
        df_tfidf = calculate_tfidf(scanner, 50)
        st.dataframe(df_tfidf, use_container_width=True)

    st.divider()
    
    # advanced sections
    
    if enable_sentiment and beta_dist:
        st.subheader("⚖️ Bayesian Sentiment Inference")
        with st.expander("🧠 How to read this chart (and why it matters)", expanded=False):
            st.markdown("""
            **The Problem:** Standard sentiment analysis gives you a single number (e.g., "52% Positive"). But is that 52% based on 5 tweets or 5 million? A single number hides that uncertainty.
            
            **The Solution:** This chart calculates the **Probability** of the true sentiment.
            *   **The Curve (PDF):** Represents likelihood. The higher the peak, the more likely that specific sentiment score is the "truth."
            *   **The Shape:** 
                *   **Narrow & Tall:** We have lots of data. We are highly confident the sentiment is exactly here.
                *   **Wide & Flat:** We don't have enough data. The true sentiment could be almost anything.
            *   **The Green Zone (95% CI):** There is a 95% probability the "True" sentiment falls within this range. 
            
            **Decision Tip:** If the green zone is very wide (e.g., spanning 30% to 70%), **do not** make business decisions based on sentiment yet; you need more data.
            """)

        bayes_result = perform_bayesian_sentiment_analysis(combined_counts, term_sentiments, pos_threshold, neg_threshold)
        if bayes_result:
            b_col1, b_col2 = st.columns([1, 2])
            with b_col1:
                st.metric("Positive Words Observed", f"{bayes_result['pos_count']:,}")
                st.metric("Negative Words Observed", f"{bayes_result['neg_count']:,}")
                st.info(f"Mean Expected Positive Rate: **{bayes_result['mean_prob']:.1%}**")
                st.success(f"95% Credible Interval:\n**{bayes_result['ci_low']:.1%} — {bayes_result['ci_high']:.1%}**")
            with b_col2:
                fig_bayes, ax_bayes = plt.subplots(figsize=(8, 4))
                ax_bayes.plot(bayes_result['x_axis'], bayes_result['pdf_y'], lw=2, color='blue', label='Posterior PDF')
                ax_bayes.fill_between(bayes_result['x_axis'], 0, bayes_result['pdf_y'], 
                                    where=(bayes_result['x_axis'] > bayes_result['ci_low']) & (bayes_result['x_axis'] < bayes_result['ci_high']),
                                    color='green', alpha=0.3, label='95% Credible Interval')
                ax_bayes.set_title("Bayesian Update of Sentiment Confidence", fontsize=10)
                ax_bayes.legend()
                ax_bayes.grid(True, alpha=0.2)
                st.pyplot(fig_bayes)
                plt.close(fig_bayes)

    show_graph = proc_conf.compute_bigrams and scanner.global_bigrams and st.checkbox("🕸️ Show Network Graph & Advanced Analytics", value=True)
    if show_graph:
        st.subheader("🔗 Network Graph")
        with st.expander("🛠️ Graph Settings & Physics", expanded=False):
            c1, c2, c3 = st.columns(3)
            min_edge_weight = c1.slider("Min Link Frequency", 2, 100, 2)
            max_nodes_graph = c1.slider("Max Nodes", 10, 200, 80)
            repulsion_val = c2.slider("Repulsion", 100, 3000, 1000)
            edge_len_val = c2.slider("Edge Length", 50, 500, 250)
            physics_enabled = c3.checkbox("Enable Physics", True)
            directed_graph = c3.checkbox("Directed Arrows", False)
            color_mode = c3.radio("Color By:", ["Community (Topic)", "Sentiment"], index=0)

        G = nx.DiGraph() if directed_graph else nx.Graph()
        filtered_bigrams = {k: v for k, v in scanner.global_bigrams.items() if v >= min_edge_weight}
        sorted_connections = sorted(filtered_bigrams.items(), key=lambda x: x[1], reverse=True)[:max_nodes_graph]
        
        if sorted_connections:
            G.add_edges_from((src, tgt, {'weight': w}) for (src, tgt), w in sorted_connections)
            try: deg_centrality = nx.degree_centrality(G)
            except: deg_centrality = {n: 1 for n in G.nodes()}
            community_map = {}
            ai_cluster_info = ""
            
            if color_mode == "Community (Topic)":
                G_undir = G.to_undirected() if directed_graph else G
                try:
                    communities = nx_comm.greedy_modularity_communities(G_undir)
                    cluster_descriptions = []
                    for group_id, community in enumerate(communities):
                        top_in_cluster = sorted(list(community), key=lambda x: combined_counts[x], reverse=True)[:5]
                        cluster_descriptions.append(f"- Cluster {group_id+1}: {', '.join(top_in_cluster)}")
                        for node in community: community_map[node] = group_id
                    ai_cluster_info = "\n".join(cluster_descriptions)
                except: pass

            community_colors = ["#FF4B4B", "#4589ff", "#ffa421", "#3cdb82", "#8b46ff", "#ff4b9f", "#00c0f2"]
            nodes, edges = [], []
            for node_id in G.nodes():
                size = 15 + (deg_centrality.get(node_id, 0) * 80)
                node_color = "#808080"
                if color_mode == "Sentiment" and enable_sentiment:
                    s = term_sentiments.get(node_id, 0)
                    if s >= pos_threshold: node_color = pos_color
                    elif s <= neg_threshold: node_color = neg_color
                elif color_mode == "Community (Topic)":
                    gid = community_map.get(node_id, 0)
                    node_color = community_colors[gid % len(community_colors)]

                # re-added font config for white, legible text
                nodes.append(Node(
                    id=node_id, 
                    label=node_id, 
                    size=size, 
                    color=node_color,
                    font={'color': 'white', 'size': 20, 'strokeWidth': 2, 'strokeColor': '#000000'}
                ))

            for (source, target), weight in sorted_connections:
                width = 1 + math.log(weight) * 0.8
                edges.append(Edge(source=source, target=target, width=width, color="#e0e0e0"))
            
            # re-added interaction dict for zoom/pan buttons
            config = Config(
                width=1000, 
                height=700, 
                directed=directed_graph, 
                physics=physics_enabled, 
                hierarchy=False, 
                interaction={"navigationButtons": True, "zoomView": True}, 
                physicsSettings={"solver": "forceAtlas2Based", "forceAtlas2Based": {"gravitationalConstant": -abs(repulsion_val), "springLength": edge_len_val, "springConstant": 0.05, "damping": 0.4}}
            )
            
            st.info("💡 **Navigation Tip:** Use the buttons in the **bottom-right** of the graph to Zoom & Pan.")
            agraph(nodes=nodes, edges=edges, config=config)
            
            # graph analytics tabs
            tab_g1, tab_g2, tab_g3, tab_g4 = st.tabs(["Basic Stats", "Top Nodes", "Text Stats", "🔥 Heatmap"])
            with tab_g1:
                col_b1, col_b2, col_b3 = st.columns(3)
                col_b1.metric("Nodes", G.number_of_nodes())
                col_b2.metric("Edges", G.number_of_edges())
                try: col_b3.metric("Density", f"{nx.density(G):.4f}")
                except: pass
            with tab_g2:
                node_weights = {n: 0 for n in G.nodes()}
                for u, v, data in G.edges(data=True):
                    w = data.get('weight', 1)
                    node_weights[u] += w
                    node_weights[v] += w
                st.dataframe(pd.DataFrame(list(node_weights.items()), columns=["Node", "Weighted Degree"]).sort_values("Weighted Degree", ascending=False).head(50), use_container_width=True)
            with tab_g3:
                 c1, c2, c3 = st.columns(3)
                 c1.metric("Total Tokens", f"{text_stats['Total Tokens']:,}")
                 c2.metric("Unique Vocab", f"{text_stats['Unique Vocabulary']:,}")
                 c3.metric("Lexical Diversity", f"{text_stats['Lexical Diversity']}")
            with tab_g4:
                # hybrid heatmap-QR generator
                viz_mode = st.radio("Visualization Mode", ["Standard Heatmap", "Hybrid Signature (Scanable)"], horizontal=True, label_visibility="collapsed")
                
                top_20 = [w for w, c in combined_counts.most_common(20)]
                
                if len(top_20) > 1:
                    # 1. Generate the Matrix Data
                    mat = np.zeros((len(top_20), len(top_20)))
                    for i, w1 in enumerate(top_20):
                        for j, w2 in enumerate(top_20):
                            if i != j: mat[i][j] = scanner.global_bigrams.get((w1, w2), 0) + scanner.global_bigrams.get((w2, w1), 0)
                    
                    # 2. Plot Heatmap to a PIL Image (Memory Buffer)
                    fig_h, ax_h = plt.subplots(figsize=(10, 10)) # Square aspect ratio is better for QR
                    ax_h.imshow(mat, cmap=colormap, interpolation='nearest') # 'nearest' gives distinct blocks
                    
                    # Clean up the plot for the "Signature" look (remove axis noise if in QR mode)
                    if viz_mode == "Hybrid Signature (Scanable)":
                        ax_h.axis('off')
                        plt.tight_layout(pad=0)
                    else:
                        ax_h.set_xticks(np.arange(len(top_20)))
                        ax_h.set_yticks(np.arange(len(top_20)))
                        ax_h.set_xticklabels(top_20, rotation=45, ha="right")
                        ax_h.set_yticklabels(top_20)
                    
                    # Save Matplotlib fig to PIL Object
                    buf = BytesIO()
                    fig_h.savefig(buf, format="png", bbox_inches="tight", pad_inches=0)
                    buf.seek(0)
                    
                    if viz_mode == "Standard Heatmap":
                        st.pyplot(fig_h)
                    
                    elif viz_mode == "Hybrid Signature (Scanable)":
                        if qrcode is None:
                            st.error("🚨 Please install: `pip install qrcode[pil]`")
                        else:
                            from PIL import Image, ImageEnhance
                            
                            # A. Prepare the Heatmap (The "Artistic Background")
                            heatmap_img = Image.open(buf).convert("RGBA")
                            
                            # Brighten heatmap so dark QR dots stand out against it
                            enhancer = ImageEnhance.Brightness(heatmap_img)
                            heatmap_img = enhancer.enhance(1.5) 
                            
                            # B. Generate the QR (The "Data Layer")
                            signature_payload = (
                                f"SIGNAL FOUNDRY\nRef: {st.session_state.get('last_sketch_hash', 'SESSION')}\n"
                                f"Top: {', '.join(top_20[:3])}"
                            )
                            qr = qrcode.QRCode(error_correction=qrcode.constants.ERROR_CORRECT_H, border=1)
                            qr.add_data(signature_payload)
                            qr.make(fit=True)
                            
                            # Create QR image: Black Data, Transparent Background
                            qr_img = qr.make_image(fill_color="black", back_color="transparent").convert("RGBA")
                            
                            # C. Composite (Merge)
                            # Resize heatmap to match QR exactly
                            heatmap_resized = heatmap_img.resize(qr_img.size)
                            
                            # Overlay: The Heatmap is the "Paper", the QR is the "Ink"
                            # We create a new image combining them
                            hybrid_img = Image.alpha_composite(heatmap_resized, qr_img)
                            
                            c1, c2 = st.columns([2, 1])
                            with c1:
                                st.image(hybrid_img, caption="Scan this Heatmap to verify analysis data.", use_container_width=True)
                            with c2:
                                st.markdown("### 🧬 Hybrid Signature")
                                st.info("The colors represent the data relationships (Heatmap). The dark overlay pattern encodes the document metadata (QR).")
                                
                                # Convert hybrid to bytes for download
                                final_buf = BytesIO()
                                hybrid_img.save(final_buf, format="PNG")
                                st.download_button("📥 Download Signature", final_buf.getvalue(), "heatmap_signature.png", "image/png")
                else:
                    st.info("Not enough data to generate signature.")
                #

    st.subheader("🔍 Bayesian Theme Discovery")
    if len(scanner.topic_docs) > 5 and DictVectorizer:
        with st.spinner(f"Running {topic_model_type} Topic Modeling..."):
            topics = perform_topic_modeling(scanner.topic_docs, n_topics, topic_model_type)
        if topics:
            cols = st.columns(len(topics))
            for idx, topic in enumerate(topics):
                with cols[idx]:
                    st.markdown(f"**Topic {topic['id']}**")
                    for w in topic['words']: st.markdown(f"`{w}`")
    else:
        st.info("Needs more data/docs to model topics.")

    # detailed frequency tables
    st.divider()
    st.subheader(f"📊 Frequency Tables (Top {top_n})")
    most_common = combined_counts.most_common(top_n)
    data = []
    if enable_sentiment:
        for w, f in most_common:
            score = term_sentiments.get(w, 0.0)
            category = get_sentiment_category(score, pos_threshold, neg_threshold)
            data.append([w, f, score, category])
    else:
        data = [[w, f] for w, f in most_common]

    cols = ["word", "count"] + (["sentiment", "category"] if enable_sentiment else [])
    st.dataframe(pd.DataFrame(data, columns=cols), use_container_width=True)
    
    if proc_conf.compute_bigrams and scanner.global_bigrams:
        st.write("Bigrams (By Frequency)")
        top_bg = scanner.global_bigrams.most_common(top_n)
        bg_data = []
        if enable_sentiment:
            for bg_tuple, f in top_bg:
                bg_str = " ".join(bg_tuple)
                score = term_sentiments.get(bg_str, 0.0)
                category = get_sentiment_category(score, pos_threshold, neg_threshold)
                bg_data.append([bg_str, f, score, category])
        else:
            bg_data = [[" ".join(bg), f] for bg, f in top_bg]
        bg_cols = ["bigram", "count"] + (["sentiment", "category"] if enable_sentiment else [])
        st.dataframe(pd.DataFrame(bg_data, columns=bg_cols), use_container_width=True)

        # NPMI in expander (original style)
        with st.expander("🔬 Phrase Significance (NPMI Score)", expanded=False):
            st.markdown("""
            **NPMI (Normalized Pointwise Mutual Information)** finds words that *belong* together, rather than just words that appear often.
            *   High Score (> 0.5): Strong association (e.g., "Artificial Intelligence").
            *   Low Score (< 0.1): Random association (e.g., "of the").
            """)
            df_npmi = calculate_npmi(scanner.global_bigrams, combined_counts, scanner.total_rows_processed)
            st.dataframe(df_npmi.head(top_n), use_container_width=True)

# --- AI analyst (restored full mode)
if combined_counts and st.session_state['authenticated']:
    st.divider()
    st.subheader("🤖 AI Analyst")
    
    top_u = [w for w, c in combined_counts.most_common(50)]
    top_b = [" ".join(bg) for bg, c in scanner.global_bigrams.most_common(20)]
    ai_ctx_str = f"Top Words: {', '.join(top_u)}\nTop Bigrams: {', '.join(top_b)}\nGraph Clusters: {locals().get('ai_cluster_info', 'N/A')}"
    
    col_ai_1, col_ai_2 = st.columns(2)
    
    with col_ai_1:
        st.markdown("**1. One-Click Theme Detection**")
        if st.button("✨ Identify Key Themes", type="primary"):
            with st.status("Analyzing..."):
                system_prompt = "You are a qualitative data analyst. Analyze the provided word frequency lists to identify 3 key themes, potential anomalies, and a summary of the subject matter."
                user_prompt = f"Data Context:\n{ai_ctx_str}"
                response = call_llm_and_track_cost(system_prompt, user_prompt, ai_config)
                st.session_state["ai_response"] = response
                st.rerun()

    with col_ai_2:
        st.markdown("**2. Ask the Data**")
        user_question = st.text_area("Ask a specific question:", height=100, placeholder="e.g., 'What are the main complaints about pricing?'")
        if st.button("Ask Question"):
            if user_question.strip():
                with st.status("Thinking..."):
                    system_prompt = "You are an expert analyst. Answer the user's question based ONLY on the provided summary statistics (word counts and associations). If you cannot answer from the data, say so."
                    user_prompt = f"Data Context:\n{ai_ctx_str}\n\nUser Question: {user_question}"
                    response = call_llm_and_track_cost(system_prompt, user_prompt, ai_config)
                    st.session_state["ai_response"] = f"**Q: {user_question}**\n\n{response}"
                    st.rerun()
            else:
                st.warning("Please enter a question.")

    if st.session_state["ai_response"]:
        st.divider()
        st.markdown("### 📋 AI Output")
        st.markdown(st.session_state["ai_response"])

st.markdown("---")
st.markdown(
    "<div style='text-align: center; color: #808080; font-size: 12px;'>"
    "Open Source software licensed under the MIT License."
    "</div>", 
    unsafe_allow_html=True
)
